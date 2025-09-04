import os
from pptx import Presentation
from PIL import Image
import pytesseract
import comtypes.client  # for slide export to images (Windows only)

# Path setup
PPT_FILE = "mydeck.pptx"
OUTPUT_DIR = "slides_output"
os.makedirs(OUTPUT_DIR, exist_ok=True)

# Step 1: Extract structured text using python-pptx
def extract_text_from_pptx(ppt_path):
    prs = Presentation(ppt_path)
    slide_texts = {}

    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        slide_texts[f"Slide_{i+1}"] = {"pptx_text": texts, "ocr_text": ""}
    return slide_texts


# Step 2: Export slides as images using PowerPoint COM API (Windows only)
def export_slides_to_images(ppt_path, output_dir):
    powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
    powerpoint.Visible = 1

    presentation = powerpoint.Presentations.Open(ppt_path, WithWindow=False)
    presentation.SaveAs(os.path.abspath(output_dir), 17)  # 17 = export to PNG
    presentation.Close()
    powerpoint.Quit()


# Step 3: OCR on slide images
def run_ocr_on_images(output_dir, slide_texts):
    for i, filename in enumerate(sorted(os.listdir(output_dir))):
        if filename.endswith(".PNG") or filename.endswith(".png"):
            img_path = os.path.join(output_dir, filename)
            img = Image.open(img_path)
            text = pytesseract.image_to_string(img)
            slide_key = f"Slide_{i+1}"
            if slide_key in slide_texts:
                slide_texts[slide_key]["ocr_text"] = text.strip()
    return slide_texts


# Step 4: Combine results
def process_ppt(ppt_file, output_dir):
    print("[INFO] Extracting PPT text...")
    slide_texts = extract_text_from_pptx(ppt_file)

    print("[INFO] Exporting slides to images...")
    export_slides_to_images(ppt_file, output_dir)

    print("[INFO] Running OCR on slide images...")
    results = run_ocr_on_images(output_dir, slide_texts)

    return results


if __name__ == "__main__":
    final_output = process_ppt(PPT_FILE, OUTPUT_DIR)

    # Print results
    for slide, content in final_output.items():
        print(f"\n--- {slide} ---")
        print("PPTX Text:")
        for t in content["pptx_text"]:
            print(f"  - {t}")
        print("OCR Text:")
        print(content["ocr_text"])
